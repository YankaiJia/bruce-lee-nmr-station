from PIL import Image
from fpdf import FPDF
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches
import os, re


def make_pdf(imagelist_here):
    pdf = FPDF()
    for image in imagelist_here:
        cover = Image.open(image)
        width, height = cover.size
        # Convert pixels to mm with 1px=0.264583 mm
        width, height = float(width * 0.264583), float(height * 0.264583)

        pdf.add_page(orientation='P', format=(width, height))
        pdf.image(image, x=0, y=0, w=width, h=height)

    return pdf


def make_slides(imagelist):
    # Create a presentation
    prs = Presentation()
    # Set slide dimensions to standard 4:3 or 16:9 if needed
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    for image_path in imagelist:
        # Create a blank slide layout
        slide_layout = prs.slide_layouts[6]  # 6 = Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Open the image to get size
        im = Image.open(image_path)
        width_px, height_px = im.size
        width_in = width_px / im.info.get('dpi', (96, 96))[0] if 'dpi' in im.info else width_px / 96
        height_in = height_px / im.info.get('dpi', (96, 96))[1] if 'dpi' in im.info else height_px / 96

        # Max dimensions for slide
        max_width = prs.slide_width
        max_height = prs.slide_height

        # Resize proportionally if needed
        aspect = width_in / height_in
        if width_in > (max_width.inches) or height_in > (max_height.inches):
            if aspect > 1:
                width = max_width
                height = Inches(max_width.inches / aspect)
            else:
                height = max_height
                width = Inches(max_height.inches * aspect)
        else:
            width = Inches(width_in)
            height = Inches(height_in)

        # Center image
        left = (prs.slide_width - width) / 2
        top = (prs.slide_height - height) / 2

        slide.shapes.add_picture(image_path, left, top, width, height)

    return prs

def get_imagelist(image_name):
    root = tk.Tk()
    root.withdraw()
    result_folder = filedialog.askdirectory(title="Select Folder with Images")
    reaction_folder_list = [os.path.join(result_folder, d) for d in os.listdir(result_folder) if os.path.isdir(os.path.join(result_folder, d))]
    reaction_folder_list = [folder for folder in reaction_folder_list if '1D EXTENDED' in folder]

    print(reaction_folder_list)
    def get_sample_index(string):
        return int(re.search(r'(\d+)-1D EXTENDED', string).group(1))
    # sort the reaction_folder_list by sample index
    # reaction_folder_list.sort(key=get_sample_index)
    print(reaction_folder_list)
    imagelist = [os.path.join(folder, image_name) for folder in reaction_folder_list]

    imagelist = [image for image in imagelist if os.path.exists(image)]

    return result_folder, imagelist

if __name__ == "__main__":

    # image_name = 'hardy_fit_diagnostic_plot.png'
    image_name = 'fitting_results.png'

    result_folder, imagelist = get_imagelist(image_name = image_name)

    pdf = make_pdf(imagelist)
    pdf.output(result_folder + f"\\all_{image_name}.pdf", "F")

    prs = make_slides(imagelist)
    prs.save(result_folder + f"\\all_{image_name}.pptx")
